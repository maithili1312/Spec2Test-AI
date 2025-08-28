import streamlit as st
import pandas as pd
import PyPDF2
from docx import Document
from pptx import Presentation
import pyshark
import re
import os
from groq import Groq  # pip install groq

import streamlit as st

# ✅ Custom Stylish UI Inject
st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Poppins:wght@400;600&display=swap');

    html, body, [class*="css"] {
        font-family: 'Poppins', sans-serif;
        background-color: #f5f7fa;
    }

    h2, h3 {
        color: #2c3e50;
        font-weight: 600;
    }

    /* File uploader */
    .stFileUploader {
        border: 2px dashed #4CAF50;
        padding: 12px;
        border-radius: 10px;
        background-color: white;
    }

    /* Buttons */
    div.stButton > button {
        background-color: #4CAF50;
        color: white;
        border: none;
        border-radius: 8px;
        font-weight: 600;
        font-size: 16px;
        padding: 8px 20px;
        transition: all 0.3s ease;
    }
    div.stButton > button:hover {
        background-color: #45a049;
        transform: scale(1.03);
    }

    /* Text input */
    input[type="text"] {
        border: 1px solid #4CAF50;
        border-radius: 8px;
        padding: 8px;
    }
    </style>
""", unsafe_allow_html=True)



# --- Streamlit Config ---
st.set_page_config(page_title="AI Test Case Generator (Groq)", page_icon="🧪")

# --- Dark Mode Toggle ---
dark_mode = st.checkbox("🌙 Enable Dark Mode")

# --- Theme Styling ---
if dark_mode:
    st.markdown("""
        <style>
        /* App background & text */
        body, .stApp {
            background: linear-gradient(135deg, #1b1b1b, #333333 100%) !important;
            color: white !important;
        }
        h1, h2, h3, label, p, span, div, .stMarkdown {
            color: white !important;
        }
        /* Dataframe dark mode */
        .stDataFrame, .dataframe {
            background-color: #222 !important;
            color: white !important;
        }
        /* Buttons */
        div.stButton > button {
            background-color: #4CAF50 !important;
            color: white !important;
            border: none;
            border-radius: 8px;
            font-weight: 600;
            font-size: 16px;
            padding: 8px 20px;
            transition: all 0.3s ease;
        }
        div.stButton > button:hover {
            background-color: #45a049 !important;
            transform: scale(1.03);
        }
        /* File uploader container */
        [data-testid="stFileUploader"] section {
            background-color: #2a2a2a !important;
            border: 2px solid #555 !important;
            border-radius: 8px;
        }
        [data-testid="stFileUploader"] section:hover {
            border-color: #4CAF50 !important;
        }
        /* File uploader label text */
        [data-testid="stFileUploader"] label {
            color: white !important;
            font-weight: 600 !important;
        }
        /* Browse files button */
        [data-testid="stFileUploader"] div div div button {
            background-color: #4CAF50 !important;
            color: white !important;
            border: none !important;
            border-radius: 6px !important;
            padding: 6px 14px !important;
            font-weight: 600 !important;
        }
        [data-testid="stFileUploader"] div div div button:hover {
            background-color: #45a049 !important;
        }
        /* Text input dark mode */
        input, textarea {
            background-color: #2a2a2a !important;
            color: white !important;
            border: 1px solid #555 !important;
            border-radius: 6px;
        }
        input::placeholder, textarea::placeholder {
            color: #bbb !important;
        }
        </style>
    """, unsafe_allow_html=True)

else:
    st.markdown("""
        <style>
        body, .stApp {
            background-color: white !important;
            color: black !important;
        }
        h1, h2, h3, label, p, span, div, .stMarkdown {
            color: black !important;
        }
        div.stButton > button {
            background-color: #4CAF50 !important;
            color: white !important;
            border: none;
            border-radius: 8px;
            font-weight: 600;
            font-size: 16px;
            padding: 8px 20px;
            transition: all 0.3s ease;
        }
        div.stButton > button:hover {
            background-color: #45a049 !important;
            transform: scale(1.03);
        }
        </style>
    """, unsafe_allow_html=True)


# --- Constants ---
MAX_FILE_SIZE = 200 * 1024 * 1024  # 200 MB



# Set your Groq API key here or in .env file
GROQ_API_KEY = os.getenv("GROQ_API_KEY", "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")
client = Groq(api_key=GROQ_API_KEY)

# --- Load File Content ---
def load_file(file):
    try:
        if file.size > MAX_FILE_SIZE:
            st.error("File exceeds 200MB limit.")
            return None

        file_type = file.type
        if file_type in ["text/plain", "application/octet-stream"]:
            return file.read().decode("utf-8")
        elif file_type == "application/vnd.tcpdump.pcap":
            capture = pyshark.FileCapture(file)
            return "\n".join([str(packet) for packet in capture])
        elif file_type == "application/pdf":
            reader = PyPDF2.PdfReader(file)
            return ''.join([page.extract_text() for page in reader.pages if page.extract_text()])
        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = Document(file)
            return '\n'.join([p.text for p in doc.paragraphs if p.text])
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            ppt = Presentation(file)
            return '\n'.join(
                [shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text") and shape.text]
            )
        else:
            st.error("Unsupported file type.")
            return None
    except Exception as e:
        st.error(f"Error reading file: {str(e)}")
        return None

# --- Ask Groq ---
def ask_groq(prompt, context=""):
    try:
        input_text = f"""
You are a QA engineer. Based on the following requirement document, generate test cases.

Context:
{context}

Instruction:
ONLY return the test cases in the following format:
<Actual test description> || <Steps to be followed> || <Expected Results>

Do NOT use labels like 'Test Case 1'. Just plain rows with || separator.
No headers, no markdown, no numbering, no explanations — only valid data rows.
"""
        response = client.chat.completions.create(
            model="llama3-70b-8192",  # Very fast & accurate Groq model
            messages=[{"role": "user", "content": input_text}],
            temperature=0.3,
            max_tokens=1500
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        st.error(f"❌ Unexpected error in ask_groq: {str(e)}")
        return None

# --- Parse and Clean Output ---
def parse_test_cases(text):
    lines = [line.strip() for line in text.split("\n") if "||" in line]
    rows = []
    for line in lines:
        if re.match(r"(?i)test case\s*\d+\s*\|\|", line):
            line = re.sub(r"(?i)test case\s*\d+\s*\|\|", "", line)

        parts = [p.strip() for p in line.split("||")]
        if len(parts) == 3:
            rows.append(parts)

    return pd.DataFrame(rows, columns=["Test Description", "Steps to be followed", "Expected Results"])

# --- Streamlit UI ---
st.markdown("## ✅ AI-Powered Test Case Generator")
st.markdown("Upload a requirements document and generate structured test cases instantly.")

file = st.file_uploader("📄 Upload your document", type=["txt", "log", "csv", "pdf", "pptx", "docx", "pcap"])
if file:
    file_content = load_file(file)
    if file_content:
        st.session_state["file_content"] = file_content
        st.success("✅ Document uploaded successfully!")

user_prompt = st.text_input("✍️ Enter your prompt")

if st.button("🚀 Generate Test Cases"):
    if user_prompt and "file_content" in st.session_state:
        with st.spinner("⚡ Generating using Groq API..."):
            response = ask_groq(user_prompt, context=st.session_state["file_content"])
            if not response:
                st.error("❌ Failed to generate test cases.")
            else:
                df = parse_test_cases(response)
                if df.empty:
                    st.warning("⚠️ No valid test cases parsed. Try a clearer prompt.")
                else:
                    st.markdown("### 📋 Test Case Table Preview")
                    st.dataframe(df, use_container_width=True)
                    csv = df.to_csv(index=False).encode("utf-8")
                    st.download_button("📥 Download as CSV", csv, "test_cases.csv", "text/csv")
    else:
        st.warning("Please upload a file and enter a prompt before generating.")
